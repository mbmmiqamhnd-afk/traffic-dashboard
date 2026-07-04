import streamlit as st
from menu import show_sidebar
show_sidebar()
import io
import os
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import white, black
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 設定網頁標題 ---
st.set_page_config(page_title="商標頁碼工具", page_icon="📝")
st.header("📝 檔案商標遮蓋與頁碼工具")

# --- 自動偵測字型 ---
def get_font_path():
    possible_paths = ["kaiu.ttf", "font.ttf", "pages/kaiu.ttf", "C:/Windows/Fonts/kaiu.ttf"]
    for p in possible_paths:
        if os.path.exists(p):
            return p
    return None

font_path = get_font_path()

# --- PDF 遮蓋邏輯（動態寬度） ---
def create_pdf_overlay(page_width, page_height, page_num, current_font):
    packet = io.BytesIO()
    c = canvas.Canvas(packet, pagesize=(page_width, page_height))

    text = f"交通組製 - 第 {page_num} 頁"
    font_size = 14

    text_width = c.stringWidth(text, current_font, font_size)
    box_width = text_width + 16
    box_height = 20

    rect_x = page_width - box_width
    rect_y = 0

    c.setFillColor(white)
    c.rect(rect_x, rect_y, box_width, box_height, fill=1, stroke=0)

    c.setFillColor(black)
    c.setFont(current_font, font_size)
    c.drawRightString(page_width - 8, 4, text)

    c.save()
    packet.seek(0)
    return packet

# --- 圖片處理邏輯 ---
def process_image(image_file, font_p):
    img = Image.open(image_file).convert("RGB")
    draw = ImageDraw.Draw(img)
    width, height = img.size

    box_w, box_h = 250, 50
    rect_x0, rect_y0 = width - box_w, height - box_h

    draw.rectangle([rect_x0, rect_y0, width, height], fill="white")

    try:
        font = ImageFont.truetype(font_p, 24) if font_p else ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()

    text = "交通組製"
    draw.text((rect_x0 + 10, rect_y0 + 10), text, fill="black", font=font)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='JPEG')
    return img_byte_arr.getvalue()

# --- PPTX 處理邏輯（動態寬度版） ---
def process_pptx(pptx_file, font_p):
    prs = Presentation(pptx_file)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    font_size_pt = 14
    padding_pt = 8  # 左右各 8pt，跟 PDF 版本一致
    box_height_pt = 22

    # 用 PIL 量測文字寬度（換算成 pt）
    def measure_text_width_pt(text, size_pt):
        try:
            # PIL 用像素單位，1pt = 96/72 px（標準換算）
            px_size = int(size_pt * 96 / 72)
            font = ImageFont.truetype(font_p, px_size) if font_p else ImageFont.load_default()
            width_px = font.getlength(text)
            return width_px * 72 / 96
        except Exception:
            # 量測失敗時的保守估計（每字約 1 個字高）
            return len(text) * size_pt

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        text = f"交通組製 - 第 {page_num} 頁"

        text_width_pt = measure_text_width_pt(text, font_size_pt)
        box_width_pt = text_width_pt + padding_pt * 2
        box_width = Pt(box_width_pt)
        box_height = Pt(box_height_pt)

        left = slide_width - box_width
        top = slide_height - box_height

        # 白色遮蓋矩形（無外框線）
        shape = slide.shapes.add_shape(1, left, top, box_width, box_height)  # 1 = MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.fill.background()
        shape.shadow.inherit = False

        # 文字
        tf = shape.text_frame
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(padding_pt)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run.font.name = "標楷體"

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()

# --- 主介面 ---
uploaded_file = st.file_uploader("上傳 PDF、PPTX 或 圖片 (JPG/PNG)", type=["pdf", "pptx", "jpg", "jpeg", "png"])

if uploaded_file and st.button("開始加工"):
    file_ext = uploaded_file.name.split('.')[-1].lower()

    try:
        # 處理 PDF
        if file_ext == "pdf":
            reader = PdfReader(uploaded_file)
            writer = PdfWriter()

            pdf_font = "Helvetica"  # 預設字型
            if font_path:
                try:
                    pdfmetrics.registerFont(TTFont('CustomFont', font_path))
                    pdf_font = 'CustomFont'
                except Exception as font_e:
                    st.warning(f"字體載入失敗，改用預設字體 (錯誤: {font_e})")

            for i, page in enumerate(reader.pages):
                w, h = float(page.mediabox.width), float(page.mediabox.height)
                overlay = create_pdf_overlay(w, h, i + 1, pdf_font)

                page.merge_page(PdfReader(overlay).pages[0])
                writer.add_page(page)

            out = io.BytesIO()
            writer.write(out)
            st.success("🎉 PDF 加工完成！")
            st.download_button("📥 下載加工版 PDF", out.getvalue(), "加工版.pdf", "application/pdf")

        # 處理 PPTX
        elif file_ext == "pptx":
            result_pptx = process_pptx(uploaded_file, font_path)
            st.success("🎉 PPTX 加工完成！")
            st.download_button(
                "📥 下載加工版 PPTX",
                result_pptx,
                "加工版.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        # 處理圖片
        else:
            result_img = process_image(uploaded_file, font_path)
            st.image(result_img, caption="預覽加工後的圖片")
            st.success("🎉 圖片加工完成！")
            st.download_button("📥 下載加工版圖片", result_img, f"processed_{uploaded_file.name}", f"image/{file_ext}")

    except Exception as e:
        st.error(f"發生錯誤: {e}")
