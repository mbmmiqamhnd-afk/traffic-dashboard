import streamlit as st
from menu import show_sidebar

# 必須在最前面呼叫
st.set_page_config(page_title="綜合檔案加工與轉檔中心", page_icon="🗂️")
show_sidebar()

import io
import os
import zipfile
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import white, black
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pdf2image import convert_from_bytes

st.header("🗂️ 綜合檔案加工與轉檔中心")
st.write("支援：檔案商標遮蓋添加頁碼、PDF 轉 PPTX/圖片、以及多圖合併轉 PDF")

# ==========================================
# 共用函式區 (商標與頁碼加工使用)
# ==========================================
def get_font_path():
    possible_paths = ["kaiu.ttf", "font.ttf", "pages/kaiu.ttf", "C:/Windows/Fonts/kaiu.ttf"]
    for p in possible_paths:
        if os.path.exists(p):
            return p
    return None

font_path = get_font_path()

def set_east_asian_font(run, font_name):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font_name)

def create_pdf_overlay(page_width, page_height, page_num, current_font):
    packet = io.BytesIO()
    c = canvas.Canvas(packet, pagesize=(page_width, page_height))

    text = f"交通組製 - 第 {page_num} 頁"
    font_size = 14

    text_width = c.stringWidth(text, current_font, font_size)
    box_width = text_width + 16
    box_height = 20

    rect_x = page_width - box_width
    rect_y = 0

    c.setFillColor(white)
    c.rect(rect_x, rect_y, box_width, box_height, fill=1, stroke=0)

    c.setFillColor(black)
    c.setFont(current_font, font_size)
    c.drawRightString(page_width - 8, 4, text)

    c.save()
    packet.seek(0)
    return packet

def process_image(image_file, font_p):
    img = Image.open(image_file).convert("RGB")
    draw = ImageDraw.Draw(img)
    width, height = img.size

    box_w, box_h = 250, 50
    rect_x0, rect_y0 = width - box_w, height - box_h

    draw.rectangle([rect_x0, rect_y0, width, height], fill="white")

    try:
        font = ImageFont.truetype(font_p, 24) if font_p else ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()

    text = "交通組製"
    draw.text((rect_x0 + 10, rect_y0 + 10), text, fill="black", font=font)

    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='JPEG')
    return img_byte_arr.getvalue()

def process_pptx(pptx_file, font_p):
    prs = Presentation(pptx_file)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    font_size_pt = 14
    padding_pt = 8  
    box_height_pt = 22

    def measure_text_width_pt(text, size_pt):
        try:
            px_size = int(size_pt * 96 / 72)
            font = ImageFont.truetype(font_p, px_size) if font_p else ImageFont.load_default()
            width_px = font.getlength(text)
            return width_px * 72 / 96
        except Exception:
            return len(text) * size_pt

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        text = f"交通組製 - 第 {page_num} 頁"

        text_width_pt = measure_text_width_pt(text, font_size_pt)
        box_width_pt = text_width_pt + padding_pt * 2
        box_width = Pt(box_width_pt)
        box_height = Pt(box_height_pt)

        left = slide_width - box_width
        top = slide_height - box_height

        shape = slide.shapes.add_shape(1, left, top, box_width, box_height) 
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(padding_pt)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run.font.name = "標楷體"
        set_east_asian_font(run, "標楷體")

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()

# ==========================================
# 介面分頁 (Tabs)
# ==========================================
tab1, tab2, tab3 = st.tabs(["📝 商標與頁碼加工", "📄 PDF 轉其他格式", "🖼️ 圖片合併轉 PDF"])

# --- Tab 1: 商標與頁碼加工 ---
with tab1:
    st.subheader("檔案商標遮蓋與頁碼添加")
    watermark_file = st.file_uploader("上傳 PDF、PPTX 或 圖片 (JPG/PNG)", type=["pdf", "pptx", "jpg", "jpeg", "png"], key="watermark")

    if watermark_file and st.button("🚀 開始加工"):
        file_ext = watermark_file.name.split('.')[-1].lower()

        try:
            if file_ext == "pdf":
                reader = PdfReader(watermark_file)
                writer = PdfWriter()

                pdf_font = "Helvetica" 
                if font_path:
                    try:
                        pdfmetrics.registerFont(TTFont('CustomFont', font_path))
                        pdf_font = 'CustomFont'
                    except Exception as font_e:
                        st.warning(f"字體載入失敗，改用預設字體 (錯誤: {font_e})")

                for i, page in enumerate(reader.pages):
                    w, h = float(page.mediabox.width), float(page.mediabox.height)
                    overlay = create_pdf_overlay(w, h, i + 1, pdf_font)
                    page.merge_page(PdfReader(overlay).pages[0])
                    writer.add_page(page)

                out = io.BytesIO()
                writer.write(out)
                st.success("🎉 PDF 加工完成！")
                st.download_button("📥 下載加工版 PDF", out.getvalue(), f"加工版_{watermark_file.name}", "application/pdf")

            elif file_ext == "pptx":
                result_pptx = process_pptx(watermark_file, font_path)
                st.success("🎉 PPTX 加工完成！")
                st.download_button("📥 下載加工版 PPTX", result_pptx, f"加工版_{watermark_file.name}", "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            else:
                result_img = process_image(watermark_file, font_path)
                st.image(result_img, caption="預覽加工後的圖片")
                st.success("🎉 圖片加工完成！")
                st.download_button("📥 下載加工版圖片", result_img, f"加工版_{watermark_file.name}", f"image/{file_ext}")

        except Exception as e:
            st.error(f"發生錯誤: {e}")

# --- Tab 2: PDF 轉其他格式 ---
with tab2:
    st.subheader("將 PDF 轉為 PPTX 或 圖片")
    pdf_convert_file = st.file_uploader("請上傳 PDF 檔案", type=["pdf"], key="pdf_convert")

    if pdf_convert_file:
        option = st.radio("請選擇轉換格式：", ("轉成 PPTX", "轉成圖片 (ZIP壓縮檔)"))

        if st.button(f"🚀 開始{option}"):
            with st.spinner("正在解析 PDF 並處理中..."):
                try:
                    file_bytes = pdf_convert_file.read()
                    images = convert_from_bytes(file_bytes, dpi=150)
                    
                    if images:
                        if option == "轉成 PPTX":
                            prs = Presentation()
                            w, h = images[0].size
                            prs.slide_width = int(w * 914400 / 150)
                            prs.slide_height = int(h * 914400 / 150)
                            
                            for img in images:
                                slide = prs.slides.add_slide(prs.slide_layouts[6])
                                img_stream = io.BytesIO()
                                img.save(img_stream, format='JPEG', quality=85)
                                img_stream.seek(0)
                                slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                            
                            pptx_out = io.BytesIO()
                            prs.save(pptx_out)
                            
                            st.success(f"✅ PPTX 轉換成功！共處理 {len(images)} 頁")
                            st.download_button("📥 點擊下載 PPTX", pptx_out.getvalue(), f"{pdf_convert_file.name.rsplit('.', 1)[0]}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                            
                        elif option == "轉成圖片 (ZIP壓縮檔)":
                            zip_buffer = io.BytesIO()
                            with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
                                for i, img in enumerate(images):
                                    img_stream = io.BytesIO()
                                    img.save(img_stream, format='JPEG', quality=100)
                                    zip_file.writestr(f"page_{i+1}.jpg", img_stream.getvalue())
                            
                            st.success(f"✅ 圖片轉換成功！共打包 {len(images)} 張圖片")
                            st.download_button("📥 點擊下載圖片 ZIP", zip_buffer.getvalue(), f"{pdf_convert_file.name.rsplit('.', 1)[0]}_images.zip", "application/zip")
                    else:
                        st.warning("⚠️ 此 PDF 沒有可提取的頁面。")
                        
                except Exception as e:
                    st.error(f"❌ 發生錯誤：{str(e)}")
                    st.info("提示：如果是部署在 Linux 環境，請確認是否已安裝 poppler-utils。")

# --- Tab 3: 圖片轉 PDF ---
with tab3:
    st.subheader("將多張圖片合併為單一 PDF")
    img_to_pdf_files = st.file_uploader("請上傳圖片 (支援多選，按上傳順序排列)", type=["jpg", "jpeg", "png", "bmp"], accept_multiple_files=True, key="img_to_pdf")

    if img_to_pdf_files:
        st.info(f"已選擇 {len(img_to_pdf_files)} 張圖片。")
        
        if st.button("🚀 開始合併為 PDF"):
            with st.spinner("正在處理圖片並建立 PDF..."):
                try:
                    image_list = []
                    for uploaded_img in img_to_pdf_files:
                        img = Image.open(uploaded_img)
                        if img.mode != 'RGB':
                            img = img.convert('RGB')
                        image_list.append(img)
                        
                    if image_list:
                        pdf_buffer = io.BytesIO()
                        first_image = image_list[0]
                        remaining_images = image_list[1:] if len(image_list) > 1 else []
                        
                        first_image.save(pdf_buffer, format="PDF", save_all=True, append_images=remaining_images, resolution=100.0)
                        
                        st.success(f"✅ PDF 建立成功！共包含 {len(image_list)} 頁")
                        st.download_button("📥 點擊下載合併後的 PDF", pdf_buffer.getvalue(), "merged_images.pdf", "application/pdf")
                except Exception as e:
                    st.error(f"❌ 發生錯誤：{str(e)}")
