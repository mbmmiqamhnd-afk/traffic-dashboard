import streamlit as st
import io
import zipfile
from pdf2image import convert_from_bytes
from pptx import Presentation

st.set_page_config(page_title="PDF 轉檔工具", page_icon="📂")
st.header("📂 PDF 格式轉換中心")

uploaded_file = st.file_uploader("上傳 PDF (建議用加工後的版本)", type=["pdf"])

if uploaded_file and st.button("開始轉換"):
    with st.spinner("轉換中... (若出現錯誤請確認 packages.txt 是否已建立)"):
        try:
            images = convert_from_bytes(uploaded_file.read(), dpi=150)
            
            # 製作 PPTX
            prs = Presentation()
            if images:
                w, h = images[0].size
                prs.slide_width = int(w * 914400 / 150)
                prs.slide_height = int(h * 914400 / 150)
                
            for img in images:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                img_stream = io.BytesIO()
                img.save(img_stream, format='JPEG')
                slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            pptx_out = io.BytesIO()
            prs.save(pptx_out)
            
            st.success(f"轉換成功！共 {len(images)} 頁")
            st.download_button("📥 下載 PPTX", pptx_out.getvalue(), "converted.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            
        except Exception as e:
            st.error(f"錯誤：{e}")
