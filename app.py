import streamlit as st
import pandas as pd
import io
import re
import gspread
import traceback
from datetime import datetime, timedelta
from pdf2image import convert_from_bytes
from pptx import Presentation

# ==========================================
# 0. 系統初始化與格式套件
# ==========================================
st.set_page_config(page_title="龍潭分局交通智慧戰情室", page_icon="🚓", layout="wide")

try:
    from gspread_formatting import *
    HAS_FORMATTING = True
except ImportError:
    HAS_FORMATTING = False

# ==========================================
# 1. 全局常數與設定區
# ==========================================
GOOGLE_SHEET_URL = "https://docs.google.com/spreadsheets/d/1HaFu5PZkFDUg7WZGV9khyQ0itdGXhXUakP4_BClFTUg/edit"

try:
    GCP_CREDS = dict(st.secrets.get("gcp_service_account", {}))
except:
    GCP_CREDS = None

# --- [重大違規常數] ---
MAJOR_UNIT_ORDER = ['科技執法', '聖亭所', '龍潭所', '中興所', '石門所', '高平所', '三和所', '警備隊', '交通分隊']
MAJOR_TARGETS = {'聖亭所': 1941, '龍潭所': 2588, '中興所': 1941, '石門所': 1479, '高平所': 1294, '三和所': 339, '交通分隊': 2526, '警備隊': 0, '科技執法': 6006}
MAJOR_FOOTNOTE = "重大交通違規指：「酒駕」、「闖紅燈」、「嚴重超速」、「逆向行駛」、「轉彎未依規定」、「蛇行、惡意逼車」及「不暫停讓行人」"

# --- [超載統計常數] ---
OVERLOAD_TARGETS = {'聖亭所': 20, '龍潭所': 27, '中興所': 20, '石門所': 16, '高平所': 14, '三和所': 8, '警備隊': 0, '交通分隊': 22}
OVERLOAD_UNIT_MAP = {'聖亭派出所': '聖亭所', '龍潭派出所': '龍潭所', '中興派出所': '中興所', '石門派出所': '石門所', '高平派出所': '高平所', '三和派出所': '三和所', '警備隊': '警備隊', '龍潭交通分隊': '交通分隊'}
OVERLOAD_UNIT_ORDER = ['聖亭所', '龍潭所', '中興所', '石門所', '高平所', '三和所', '警備隊', '交通分隊']

# --- [強化專案常數] ---
PROJECT_NAME = "強化交通安全執法專案勤務取締件數統計表"
PROJECT_TARGETS = {
    '聖亭所': [5, 115, 5, 16, 7, 10], '龍潭所': [6, 145, 7, 20, 9, 12],
    '中興所': [5, 115, 5, 16, 7, 10], '石門所': [3, 80, 4, 11, 5, 7],
    '高平所': [3, 80, 4, 11, 5, 7], '三和所': [2, 40, 2, 6, 2, 5],
    '交通分隊': [5, 115, 4, 16, 6, 8], '交通組': [0, 0, 0, 0, 0, 0], '警備隊': [0, 0, 0, 0, 0, 0]
}
PROJECT_CATS = ["酒後駕車", "闖紅燈", "嚴重超速", "車不讓人", "行人違規", "大型車違規"]
PROJECT_LAW_MAP = {
    "酒後駕車": ["35條", "73條2項", "73條3項"], 
    "闖紅燈": ["53條"], 
    "嚴重超速": ["43條", "40條"],  
    "車不讓人": ["44條", "48條"], 
    "行人違規": ["78條"]
}

# ==========================================
# 2. 輔助工具區
# ==========================================
def get_gsheet_rich_text_req(sheet_id, row_idx, col_idx, text):
    text = str(text)
    pattern = r'([0-9\(\)\/\-]+)'
    tokens = re.split(pattern, text)
    runs = []
    current_pos = 0
    for token in tokens:
        if not token: continue
        color = {"red": 1.0, "green": 0.0, "blue": 0.0} if re.match(pattern, token) else {"red": 0.0, "green": 0.0, "blue": 0.0}
        runs.append({"startIndex": current_pos, "format": {"foregroundColor": color, "bold": True}})
        current_pos += len(token)
    return {
        "updateCells": {
            "rows": [{"values": [{"userEnteredValue": {"stringValue": text}, "textFormatRuns": runs}]}],
            "fields": "userEnteredValue,textFormatRuns",
            "range": {"sheetId": sheet_id, "startRowIndex": row_idx, "endRowIndex": row_idx + 1, "startColumnIndex": col_idx, "endColumnIndex": col_idx + 1}
        }
    }

# ==========================================
# 3. 業務邏輯處理區
# ==========================================

def process_tech_enforcement(files):
    f = files[0]
    f.seek(0)
    df = pd.read_csv(f, encoding='cp950') if f.name.endswith('.csv') else pd.read_excel(f)
    df.columns = [str(c).strip() for c in df.columns]
    loc_col = next((c for c in df.columns if c in ['違規地點', '路口名稱', '地點']), None)
    if not loc_col: return
    df[loc_col] = df[loc_col].astype(str).str.replace('桃園市', '').str.replace('龍潭區', '').str.strip()
    yesterday = datetime.now() - timedelta(days=1)
    date_range_str = f"{yesterday.year - 1911}年1月1日至{yesterday.year - 1911}年{yesterday.month}月{yesterday.day}日"
    loc_summary = df[loc_col].value_counts().head(10).reset_index()
    loc_summary.columns = ['路段名稱', '舉發件數']
    st.write("📊 **科技執法路段排行：**")
    st.dataframe(loc_summary, hide_index=True)
    if GCP_CREDS:
        gc = gspread.service_account_from_dict(GCP_CREDS)
        sh = gc.open_by_url(GOOGLE_SHEET_URL)
        ws_name = "科技執法-路段排行"
        ws = sh.worksheet(ws_name) if ws_name in [s.title for s in sh.worksheets()] else sh.add_worksheet(title=ws_name, rows="100", cols="20")
        ws.clear()
        title_text = f"科技執法成效 ({date_range_str})"
        ws.update(range_name='A1', values=[[title_text, ""], ["路段名稱", "舉發件數"]] + loc_summary.values.tolist() + [["舉發總數", len(df)]])
        reqs = {"requests": [{"updateCells": {"range": {"sheetId": ws.id, "startRowIndex": 0, "endRowIndex": 1, "startColumnIndex": 0, "endColumnIndex": 1},
                "rows": [{"values": [{"userEnteredValue": {"stringValue": title_text},
                "textFormatRuns": [{"startIndex": 0, "format": {"foregroundColor": {"red": 0.0, "green": 0.0, "blue": 1.0}, "bold": True, "fontSize": 24}},
                                   {"startIndex": len("科技執法成效 "), "format": {"foregroundColor": {"red": 1.0, "green": 0.0, "blue": 0.0}, "bold": True, "fontSize": 24}}]}]}],
                "fields": "userEnteredValue,textFormatRuns"}}]}
        sh.batch_update(reqs)

def process_overload(files):
    f_wk, f_yt, f_ly = None, None, None
    for f in files:
        if "(1)" in f.name: f_yt = f
        elif "(2)" in f.name: f_ly = f
        else: f_wk = f
    def parse_rpt(f):
        if not f: return {}, "0000000", "0000000"
        f.seek(0)
        counts, s, e = {}, "0000000", "0000000"
        text_block = pd.read_excel(f, header=None, nrows=15).to_string()
        m = re.search(r'(\d{3,7}).*至\s*(\d{3,7})', text_block)
        if m: s, e = m.group(1), m.group(2)
        f.seek(0); xls = pd.ExcelFile(f)
        for sn in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sn, header=None)
            u = None
            for _, r in df.iterrows():
                rs = " ".join([str(x) for x in r.values])
                if "舉發單位：" in rs:
                    m2 = re.search(r"舉發單位：(\S+)", rs)
                    if m2: u = m2.group(1).strip()
                if "總計" in rs and u:
                    nums = [float(str(x).replace(',','')) for x in r if str(x).replace('.','',1).isdigit()]
                    if nums:
                        short = OVERLOAD_UNIT_MAP.get(u, u)
                        if short in OVERLOAD_UNIT_ORDER: counts[short] = counts.get(short, 0) + int(nums[-1])
                        u = None
        return counts, s, e
    d_wk, s_wk, e_wk = parse_rpt(f_wk); d_yt, s_yt, e_yt = parse_rpt(f_yt); d_ly, s_ly, e_ly = parse_rpt(f_ly)
    raw_wk, raw_yt, raw_ly = f"本期 ({s_wk[-4:]}~{e_wk[-4:]})", f"本年累計 ({s_yt[-4:]}~{e_yt[-4:]})", f"去年累計 ({s_ly[-4:]}~{e_ly[-4:]})"
    body = []
    for u in OVERLOAD_UNIT_ORDER:
        yv, tv = d_yt.get(u, 0), OVERLOAD_TARGETS.get(u, 0)
        body.append({'統計期間': u, raw_wk: d_wk.get(u, 0), raw_yt: yv, raw_ly: d_ly.get(u, 0), '本年與去年同期比較': yv - d_ly.get(u, 0), '目標值': tv, '達成率': f"{yv/tv:.0%}" if tv > 0 else "—"})
    df_body = pd.DataFrame(body)
    sum_v = df_body[df_body['統計期間'] != '警備隊'][[raw_wk, raw_yt, raw_ly, '目標值']].sum()
    total_row = pd.DataFrame([{'統計期間': '合計', raw_wk: sum_v[raw_wk], raw_yt: sum_v[raw_yt], raw_ly: sum_v[raw_ly], '本年與去年同期比較': sum_v[raw_yt] - sum_v[raw_ly], '目標值': sum_v['目標值'], '達成率': f"{sum_v[raw_yt]/sum_v['目標值']:.0%}" if sum_v['目標值'] > 0 else "0%"}])
    df_final = pd.concat([total_row, df_body], ignore_index=True)
    st.write("📊 **超載統計結果：**")
    st.dataframe(df_final, hide_index=True)
    if GCP_CREDS:
        gc = gspread.service_account_from_dict(GCP_CREDS)
        sh = gc.open_by_url(GOOGLE_SHEET_URL); ws = sh.get_worksheet(1)
        ws.update(range_name='A1', values=[['取締超載違規件數統計表']])
        ws.update(range_name='A2', values=[df_final.columns.tolist()] + df_final.values.tolist())
        requests = []
        for i, col_name in enumerate(df_final.columns):
            if "(" in col_name:
                p_start = col_name.find("(")
                requests.append({"updateCells": {"range": {"sheetId": ws.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": i, "endColumnIndex": i+1},
                        "rows": [{ "values": [{ "textFormatRuns": [{"startIndex": 0, "format": {"foregroundColor": {"red": 0.0, "green": 0.0, "blue": 0.0}, "bold": True}}, {"startIndex": p_start, "format": {"foregroundColor": {"red": 1.0, "green": 0.0, "blue": 0.0}, "bold": True}}], "userEnteredValue": {"stringValue": col_name} }] }], "fields": "userEnteredValue,textFormatRuns"}})
        if requests: sh.batch_update({"requests": requests})

# ----------------- [3. 重大交通違規] -----------------
def process_major(files):
    if len(files) < 2:
        st.error("❌ 請上傳『本期』與『年累計』報表。若要精確比較細項，請一併上傳第三份『去年累計』報表。")
        return
    f_wk, f_year, f_ly = None, None, None
    for f in files:
        if "本期" in f.name: f_wk = f
        elif "去年" in f.name: f_ly = f
        elif "年累計" in f.name: f_year = f
    if not f_wk or not f_year:
        sorted_files = sorted([f for f in files if "重大" in f.name or "重點" in f.name] or files, key=lambda x: x.size)
        if len(sorted_files) >= 1 and not f_wk: f_wk = sorted_files[0]
        if len(sorted_files) >= 2 and not f_year: f_year = sorted_files[1]
        if len(sorted_files) >= 3 and not f_ly: f_ly = sorted_files[2]

    def get_robust_date(df):
        try:
            raw_cells = [str(val) for val in df.head(10).values.flatten() if pd.notna(val)]
            clean_text = re.sub(r'\s+', '', "".join(raw_cells))
            match = re.search(r'1\d{2}(\d{4})[至\-~]1\d{2}(\d{4})', clean_text)
            if match: return match.group(1)+"-"+match.group(2)
            return ""
        except: return ""

    def clean_unit(n):
        if pd.isna(n): return None
        n = str(n).strip()
        if '分隊' in n: return '交通分隊'
        if any(k in n for k in ['科技', '交通組']): return '科技執法'
        if '警備' in n: return '警備隊'
        for k in ['聖亭', '龍潭', '中興', '石門', '高平', '三和']:
            if k in n: return k + '所'
        return None

    def to_i(v):
        try: return int(float(str(v).replace(',', '').strip()))
        except: return 0

    def get_dfs(f):
        if not f: return []
        f.seek(0)
        if f.name.lower().endswith('.csv'):
            try: return [pd.read_csv(f, header=None)]
            except: f.seek(0); return [pd.read_csv(f, encoding='cp950', header=None)]
        else:
            try: xl = pd.ExcelFile(f); return [pd.read_excel(xl, sheet_name=sn, header=None) for sn in xl.sheet_names]
            except: return []

    dfs_wk, dfs_yr, dfs_ly = get_dfs(f_wk), get_dfs(f_year), get_dfs(f_ly)

    def parse_main_table(dfs):
        d_yt, d_ly = {}, {}
        dt_str = ""
        for df in dfs:
            if not dt_str: dt_str = get_robust_date(df)
            for _, r in df.iterrows():
                u = clean_unit(r.iloc[0])
                if u and "合計" not in str(r.iloc[0]):
                    if len(r) > 16: d_yt[u] = {'stop': to_i(r.iloc[15]), 'cit': to_i(r.iloc[16])}
                    if len(r) > 19: d_ly[u] = {'stop': to_i(r.iloc[18]), 'cit': to_i(r.iloc[19])}
        return d_yt, d_ly, dt_str

    d_wk_yt, _, date_wk = parse_main_table(dfs_wk)
    d_yr_yt, d_yr_ly_internal, date_yr = parse_main_table(dfs_yr)
    d_ly_yt, _, date_ly = parse_main_table(dfs_ly)

    table_rows = []
    summary = {k: 0 for k in ['ws', 'wc', 'ys', 'yc', 'ls', 'lc', 'diff', 'tgt']}
    for u in MAJOR_UNIT_ORDER:
        w_data = d_wk_yt.get(u, {'stop':0, 'cit':0})
        y_data = d_yr_yt.get(u, {'stop':0, 'cit':0})
        l_data = d_ly_yt.get(u, {'stop':0, 'cit':0}) if dfs_ly else d_yr_ly_internal.get(u, {'stop':0, 'cit':0})
        y_total = y_data['stop'] + y_data['cit']
        l_total = l_data['stop'] + l_data['cit']
        tgt = MAJOR_TARGETS.get(u, 0); diff = int(y_total - l_total)
        rate = f"{(y_total/tgt):.1%}" if tgt > 0 else "0%"
        if u != '警備隊': summary['diff'] += diff; summary['tgt'] += tgt
        table_rows.append([u, w_data['stop'], w_data['cit'], y_data['stop'], y_data['cit'], l_data['stop'], l_data['cit'], diff if u != '警備隊' else "—", tgt, rate if u != '警備隊' else "—"])
        summary['ws'] += w_data['stop']; summary['wc'] += w_data['cit']
        summary['ys'] += y_data['stop']; summary['yc'] += y_data['cit']
        summary['ls'] += l_data['stop']; summary['lc'] += l_data['cit']
    total_rate = f"{((summary['ys']+summary['yc'])/summary['tgt']):.1%}" if summary['tgt'] > 0 else "0%"
    table_rows.insert(0, ['合計', summary['ws'], summary['wc'], summary['ys'], summary['yc'], summary['ls'], summary['lc'], summary['diff'], summary['tgt'], total_rate])
    table_rows.append([MAJOR_FOOTNOTE] + [""] * 9)
    h_wk, h_yr = f"本期({date_wk})", f"本年累計({date_yr})"
    h_ls = f"去年累計({date_ly if dfs_ly else date_yr})"
    header_1 = ['統計期間', h_wk, h_wk, h_yr, h_yr, h_ls, h_ls, '本年與去年同期比較', '目標值', '達成率']
    header_2 = ['取締方式', '當場攔停', '逕行舉發', '當場攔停', '逕行舉發', '當場攔停', '逕行舉發', '', '', '']
    df_result = pd.DataFrame(table_rows, columns=pd.MultiIndex.from_arrays([header_1, header_2]))
    st.write("📊 **重大違規統計結果 (總表)：**"); st.dataframe(df_result, use_container_width=True)

    DETAIL_CATEGORIES = {"酒駕": ["酒駕", "酒後", "35條"], "闖紅燈": ["闖紅燈", "53條"], "逆向行駛": ["逆向", "45條"], "轉彎未依規定": ["轉彎", "48條"], "蛇行惡意逼車": ["蛇行", "逼車", "惡意", "43條"], "不暫停讓行人": ["行人", "車不讓人", "暫停讓", "44條"]}
    def parse_detail_data(dfs):
        res = {cat: {u: {'stop':0, 'cit':0} for u in MAJOR_UNIT_ORDER} for cat in DETAIL_CATEGORIES}
        if not dfs: return res
        for df in dfs:
            header_idx = -1
            for i in range(min(15, len(df))):
                row_str = "".join([str(x) for x in df.iloc[i].values if pd.notna(x)])
                if sum(1 for kw in ["酒駕", "闖紅燈", "逆向行駛", "轉彎"] if kw in row_str) >= 2: header_idx = i; break
            if header_idx != -1:
                headers = [str(x).replace('\n', '').strip() for x in df.iloc[header_idx].values]
                sub_headers = [str(x).replace('\n', '').strip() for x in df.iloc[header_idx+1].values] if header_idx+1 < len(df) else headers
                cat_cols = {cat: {'stop': -1, 'cit': -1} for cat in DETAIL_CATEGORIES}
                for c in range(len(headers)):
                    h1, h2 = headers[c], sub_headers[c]
                    current_cat = next((cat for cat, kws in DETAIL_CATEGORIES.items() if any(kw in h1 for kw in kws)), None)
                    if current_cat:
                        if any(k in h2 for k in ["現場", "攔停", "當場", "違法"]):
                            if cat_cols[current_cat]['stop'] == -1: cat_cols[current_cat]['stop'] = c
                        elif any(k in h2 for k in ["逕", "違規"]):
                            if cat_cols[current_cat]['cit'] == -1: cat_cols[current_cat]['cit'] = c
                for _, row in df.iloc[header_idx+1:].iterrows():
                    u = clean_unit(row.values[0])
                    if u and "合計" not in str(row.values[0]):
                        for cat in DETAIL_CATEGORIES:
                            cs, cc = cat_cols[cat]['stop'], cat_cols[cat]['cit']
                            if cs != -1 and cs < len(row): res[cat][u]['stop'] += to_i(row.values[cs])
                            if cc != -1 and cc < len(row): res[cat][u]['cit'] += to_i(row.values[cc])
        return res

    d_yr_cat, d_ly_cat = parse_detail_data(dfs_yr), parse_detail_data(dfs_ly)
    cat_dfs = {}
    h1_cat = ['統計期間', '今年累計', '今年累計', '今年累計', '去年累計', '去年累計', '去年累計', '今年與去年同期比較', '今年與去年同期比較', '今年與去年同期比較']
    h2_cat = ['單位', '當場攔停', '逕行舉發', '合計', '當場攔停', '逕行舉發', '合計', '當場攔停', '逕行舉發', '合計']
    for cat in DETAIL_CATEGORIES.keys():
        rows = []; sum_cat = {'ys':0, 'yc':0, 'yt':0, 'ls':0, 'lc':0, 'lt':0, 'ds':0, 'dc':0, 'dt':0}
        for u in MAJOR_UNIT_ORDER:
            ys, yc = d_yr_cat[cat][u]['stop'], d_yr_cat[cat][u]['cit']
            ls, lc = (d_ly_cat[cat][u]['stop'], d_ly_cat[cat][u]['cit']) if dfs_ly else (0, 0)
            yt, lt = ys + yc, ls + lc
            ds, dc, dt = ys - ls, yc - lc, yt - lt
            rows.append([u, ys, yc, yt, ls, lc, lt, ds if u != '警備隊' else "—", dc if u != '警備隊' else "—", dt if u != '警備隊' else "—"])
            sum_cat['ys'] += ys; sum_cat['yc'] += yc; sum_cat['yt'] += yt
            sum_cat['ls'] += ls; sum_cat['lc'] += lc; sum_cat['lt'] += lt
            if u != '警備隊': sum_cat['ds'] += ds; sum_cat['dc'] += dc; sum_cat['dt'] += dt
        rows.insert(0, ['合計', sum_cat['ys'], sum_cat['yc'], sum_cat['yt'], sum_cat['ls'], sum_cat['lc'], sum_cat['lt'], sum_cat['ds'], sum_cat['dc'], sum_cat['dt']])
        cat_dfs[cat] = pd.DataFrame(rows, columns=pd.MultiIndex.from_arrays([h1_cat, h2_cat]))
    with st.expander("🔍 檢視 6 大項重大違規細表"):
        for cat, df_c in cat_dfs.items(): st.write(f"**【{cat}】統計表**"); st.dataframe(df_c, use_container_width=True)

    if GCP_CREDS:
        try:
            gc = gspread.service_account_from_dict(GCP_CREDS); sh = gc.open_by_url(GOOGLE_SHEET_URL)
            existing_sheets = [s.title for s in sh.worksheets()]; requests = []
            
            # 定義顏色
            red_color = {"red": 1.0, "green": 0.0, "blue": 0.0}
            black_color = {"red": 0.0, "green": 0.0, "blue": 0.0}
            blue_color = {"red": 0.0, "green": 0.0, "blue": 1.0} # 用於細項大標題的藍色

            # --- 3-1. 同步總表 ---
            ws_main = sh.get_worksheet(0); data_body_m = df_result.values.tolist()
            ws_main.update(range_name='A2', values=[[t[0] for t in df_result.columns], [t[1] for t in df_result.columns]] + data_body_m)
            for i, text in enumerate([t[0] for t in df_result.columns]):
                if "(" in text:
                    p_start = text.find("("); requests.append({"updateCells": {"range": {"sheetId": ws_main.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": i, "endColumnIndex": i+1}, "rows": [{ "values": [{ "textFormatRuns": [{"startIndex": 0, "format": {"foregroundColor": black_color, "bold": True}}, {"startIndex": p_start, "format": {"foregroundColor": red_color, "bold": True}}], "userEnteredValue": {"stringValue": text} }] }], "fields": "userEnteredValue,textFormatRuns"}})
            for r_idx, row_vals in enumerate(data_body_m):
                val = row_vals[7]
                if isinstance(val, (int, float)) and val < 0: requests.append({"repeatCell": {"range": {"sheetId": ws_main.id, "startRowIndex": 3+r_idx, "endRowIndex": 4+r_idx, "startColumnIndex": 7, "endColumnIndex": 8}, "cell": {"userEnteredFormat": {"textFormat": {"foregroundColor": red_color}}}, "fields": "userEnteredFormat.textFormat.foregroundColor"}})

            # --- 3-2. 同步細項表 ---
            for cat, df_c in cat_dfs.items():
                ws_name = f"重大違規-{cat}"
                ws_cat = sh.worksheet(ws_name) if ws_name in existing_sheets else sh.add_worksheet(title=ws_name, rows="30", cols="15")
                ws_cat.clear()
                title_text = f"取締【{cat}】違規統計表 (累計至 {date_yr})"
                ws_cat.update(range_name='A1', values=[[title_text]] + [[t[0] for t in df_c.columns], [t[1] for t in df_c.columns]] + df_c.values.tolist())
                
                # 細項表大標題 (A1)：藍色 + 紅色
                if "(" in title_text:
                    p_start_title = title_text.find("(")
                    requests.append({
                        "updateCells": {
                            "range": {"sheetId": ws_cat.id, "startRowIndex": 0, "endRowIndex": 1, "startColumnIndex": 0, "endColumnIndex": 1},
                            "rows": [{
                                "values": [{
                                    "userEnteredValue": {"stringValue": title_text},
                                    "textFormatRuns": [
                                        {"startIndex": 0, "format": {"foregroundColor": blue_color}},
                                        {"startIndex": p_start_title, "format": {"foregroundColor": red_color}}
                                    ]
                                }]
                            }],
                            "fields": "userEnteredValue,textFormatRuns"
                        }
                    })

                # 細項表表頭雙色 (黑色 + 紅色)
                for i, text in enumerate([t[0] for t in df_c.columns]):
                    if "(" in text:
                        p_start = text.find("(")
                        requests.append({
                            "updateCells": {
                                "range": {"sheetId": ws_cat.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": i, "endColumnIndex": i+1},
                                "rows": [{ "values": [{ "textFormatRuns": [
                                    {"startIndex": 0, "format": {"foregroundColor": black_color}},
                                    {"startIndex": p_start, "format": {"foregroundColor": red_color}}
                                ], "userEnteredValue": {"stringValue": text} }] }],
                                "fields": "userEnteredValue,textFormatRuns"
                            }
                        })

                # 合併與置中
                requests.extend([
                    {"mergeCells": {"range": {"sheetId": ws_cat.id, "startRowIndex": 0, "endRowIndex": 1, "startColumnIndex": 0, "endColumnIndex": 10}, "mergeType": "MERGE_ALL"}},
                    {"mergeCells": {"range": {"sheetId": ws_cat.id, "startRowIndex": 1, "endRowIndex": 3, "startColumnIndex": 0, "endColumnIndex": 1}, "mergeType": "MERGE_ALL"}},
                    {"mergeCells": {"range": {"sheetId": ws_cat.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": 1, "endColumnIndex": 4}, "mergeType": "MERGE_ALL"}},
                    {"mergeCells": {"range": {"sheetId": ws_cat.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": 4, "endColumnIndex": 7}, "mergeType": "MERGE_ALL"}},
                    {"mergeCells": {"range": {"sheetId": ws_cat.id, "startRowIndex": 1, "endRowIndex": 2, "startColumnIndex": 7, "endColumnIndex": 10}, "mergeType": "MERGE_ALL"}},
                    {"repeatCell": {"range": {"sheetId": ws_cat.id, "startRowIndex": 0, "endRowIndex": 3, "startColumnIndex": 0, "endColumnIndex": 10}, "cell": {"userEnteredFormat": {"horizontalAlignment": "CENTER", "verticalAlignment": "MIDDLE"}}, "fields": "userEnteredFormat.horizontalAlignment,userEnteredFormat.verticalAlignment"}}
                ])

                # 負數差異標示紅字
                for r_idx, row_vals in enumerate(df_c.values.tolist()):
                    target_row = 3 + r_idx
                    for c_idx in [7, 8, 9]:
                        val = row_vals[c_idx]
                        if isinstance(val, (int, float)) and val < 0:
                            requests.append({"repeatCell": {"range": {"sheetId": ws_cat.id, "startRowIndex": target_row, "endRowIndex": target_row + 1, "startColumnIndex": c_idx, "endColumnIndex": c_idx + 1}, "cell": {"userEnteredFormat": {"textFormat": {"foregroundColor": red_color}}}, "fields": "userEnteredFormat.textFormat.foregroundColor"}})

            if requests: sh.batch_update({"requests": requests})
            st.write("✅ 重大違規雲端同步完成！")
        except Exception as e: st.error(f"雲端同步出錯：{e}")

# ----------------- [4. 強化專案/交通事故/靜桃計畫] (保留原邏輯) -----------------
def process_project(files):
    f1 = next((f for f in files if any(k in f.name for k in ["強化", "法條", "自選匯出"])), None)
    f2_list = [f for f in files if any(k in f.name.upper() for k in ["R17", "砂石", "大貨"])]
    if not f1 or not f2_list: return
    def s_read(f, **kwargs):
        f.seek(0)
        if f.name.endswith('.csv'):
            try: return pd.read_csv(f, **kwargs)
            except: f.seek(0); return pd.read_csv(f, encoding='cp950', **kwargs)
        return pd.read_excel(f, **kwargs)
    date_str = "未知期間"
    df1_h = s_read(f1, nrows=10, header=None)
    for _, r in df1_h.iterrows():
        for c in r.values:
            if '統計期間' in str(c):
                m = re.search(r'([0-9年月日\-至\s]+)', str(c).replace('(入案日)', '').split('：')[-1].split(':')[-1].strip())
                if m: date_str = m.group(1).replace('115', '').strip()
    df1 = s_read(f1, skiprows=3).reset_index(drop=True)
    df2_all = []
    for f in f2_list:
        df_t = s_read(f, header=None)
        h_idx = next((i for i, r in df_t.head(30).iterrows() if '單位' in [str(x).strip() for x in r.values] and '舉發總數' in [str(x).strip() for x in r.values]), None)
        if h_idx is not None:
            df_c = df_t.iloc[h_idx+1:].copy()
            df_c.columns = [str(x).strip() for x in df_t.iloc[h_idx].values]
            df_2_clean = df_c.loc[:, ~df_c.columns.duplicated()]
            df2_all.append(df_2_clean)
    df2 = pd.concat(df2_all, ignore_index=True)
    for c in ['舉發總數', '違反管制規定', '其他違規']: df2[c] = pd.to_numeric(df2.get(c, 0), errors='coerce').fillna(0)
    df2['大型車純違規'] = (df2['舉發總數'] - df2['違反管制規定'] - df2['其他違規']).clip(lower=0)
    def get_unit(raw):
        raw = str(raw).strip()
        if '交通分隊' in raw: return '交通分隊'
        if '交通組' in raw: return '交通組'
        if '警備隊' in raw: return '警備隊'
        for k in ['聖亭', '中興', '石門', '高平', '三和']: 
            if k in raw: return k + '所'
        if '龍潭派出所' in raw or raw in ['龍潭', '龍潭所']: return '龍潭所'
        return None
    def get_c(unit):
        r = df1[df1.get('單位', pd.Series()).apply(get_unit) == unit]
        return {cat: int(r[[c for c in df1.columns if any(k in str(c) for k in PROJECT_LAW_MAP.get(cat, []))]].sum().sum()) if not r.empty else 0 for cat in PROJECT_CATS[:5]}
    final_rows = []
    for u, tgts in PROJECT_TARGETS.items():
        d15 = get_c(u); u_r = df2[df2['單位'].apply(get_unit) == u]
        h_sum = int(u_r['大型車純違規'].sum()) if not u_r.empty else 0
        res = [u]
        for i, cat in enumerate(PROJECT_CATS):
            cnt = d15.get(cat, 0) if cat != "大型車違規" else h_sum
            res.extend([cnt, tgts[i], f"{(cnt/tgts[i]*100):.1f}%" if tgts[i] > 0 else "0.0%"])
        final_rows.append(res)
    headers = ["單位"] + [f"{cat}_{x}" for cat in PROJECT_CATS for x in ["取締件數", "目標值", "達成率"]]
    df_f = pd.DataFrame(final_rows, columns=headers)
    t_row = ["合計"]
    for i in range(1, len(headers), 3):
        cs, ts = df_f.iloc[:, i].sum(), df_f.iloc[:, i+1].sum()
        t_row.extend([int(cs), int(ts), f"{(cs/ts*100):.1f}%" if ts > 0 else "0.0%"])
    df_f = pd.concat([pd.DataFrame([t_row], columns=headers), df_f], ignore_index=True)
    st.write(f"📊 **{PROJECT_NAME} 統計結果：**"); st.dataframe(df_f, hide_index=True)
    if GCP_CREDS:
        gc = gspread.service_account_from_dict(GCP_CREDS); ws = gc.open_by_url(GOOGLE_SHEET_URL).worksheet(PROJECT_NAME)
        full_t = f"{PROJECT_NAME} (統計期間：{date_str})"
        ws.clear(); ws.update(range_name='A1', values=[[full_t] + [""] * 18, [""] + [c for c in PROJECT_CATS for _ in range(3)], ["單位"] + ["取締件數", "目標值", "達成率"] * 6] + df_f.values.tolist())

def process_accident(files):
    meta = []
    for f in files:
        f.seek(0); df_raw = pd.read_csv(f, header=None) if f.name.endswith('.csv') else pd.read_excel(f, header=None)
        dates = re.findall(r'(\d{3})[./](\d{1,2})[./](\d{1,2})', str(df_raw.iloc[:5, :5].values))
        if len(dates) >= 2:
            df_raw[0] = df_raw[0].astype(str); df_data = df_raw[df_raw[0].str.contains("所|總計|合計", na=False)].rename(columns={0: "Station", 5: "A1_Deaths", 9: "A2_Injuries"})
            for c in ["A1_Deaths", "A2_Injuries"]: df_data[c] = pd.to_numeric(df_data[c].astype(str).str.replace(",", ""), errors='coerce').fillna(0)
            df_data['Station_Short'] = df_data['Station'].str.replace('派出所', '所').str.replace('總計', '合計').str.strip()
            meta.append({'df': df_data, 'year': int(dates[1][0]), 'start_day': int(dates[0][1])*100 + int(dates[0][2]), 'range': f"{int(dates[0][1]):02d}{int(dates[0][2]):02d}-{int(dates[1][1]):02d}{int(dates[1][2]):02d}", 'is_cumu': (int(dates[0][1]) == 1 and int(dates[0][2]) == 1)})
    this_year = max(m['year'] for m in meta); f_lst = sorted([f for f in meta if f['year'] < this_year], key=lambda x: x['year'])[-1]
    f_cur = next(f for f in meta if f['year'] == this_year and f['is_cumu'])
    period_files = sorted([f for f in meta if f['year'] == this_year and not f['is_cumu']], key=lambda x: x['start_day'])
    f_prev, f_wk = period_files[0], period_files[1]
    def bld_tbl(c_name, is_a2=False):
        m = pd.merge(f_wk['df'][['Station_Short', c_name]], f_prev['df'][['Station_Short', c_name]], on='Station_Short', suffixes=('_wk', '_prev'))
        m = pd.merge(pd.merge(m, f_cur['df'][['Station_Short', c_name]].rename(columns={c_name: c_name+'_cur'}), on='Station_Short'), f_lst['df'][['Station_Short', c_name]].rename(columns={c_name: c_name+'_lst'}), on='Station_Short')
        m = m[m['Station_Short'].isin(['聖亭所', '龍潭所', '中興所', '石門所', '高平所', '三和所'])].copy()
        m = pd.concat([pd.DataFrame([dict(m.select_dtypes(include='number').sum().to_dict(), Station_Short='合計')]), m], ignore_index=True)
        m['Diff'] = m[c_name+'_cur'] - m[c_name+'_lst']
        if is_a2:
            m['Pct'] = m.apply(lambda x: f"{(x['Diff']/x[c_name+'_lst']):.2%}" if x[c_name+'_lst'] != 0 else "0.00%", axis=1)
            res = m[['Station_Short', c_name+'_wk', c_name+'_prev', c_name+'_cur', c_name+'_lst', 'Diff', 'Pct']]
            res.columns = ['統計期間', f'本期({f_wk["range"]})', f'前期({f_prev["range"]})', f'本年累計({f_cur["range"]})', f'去年累計({f_lst["range"]})', '本年與去年同期比較', '增減比例']
        else:
            res = m[['Station_Short', c_name+'_wk', c_name+'_cur', c_name+'_lst', 'Diff']]
            res.columns = ['統計期間', f'本期({f_wk["range"]})', f'本年累計({f_cur["range"]})', f'去年累計({f_lst["range"]})', '本年與去年同期比較']
        return res
    a1_res, a2_res = bld_tbl('A1_Deaths'), bld_tbl('A2_Injuries', True)
    c1, c2 = st.columns(2); c1.write("📊 **A1 死亡人數統計**"); c1.dataframe(a1_res, hide_index=True); c2.write("📊 **A2 受傷人數統計**"); c2.dataframe(a2_res, hide_index=True)
    if GCP_CREDS:
        gc = gspread.service_account_from_dict(GCP_CREDS); sh = gc.open_by_url(GOOGLE_SHEET_URL)
        for ws_idx, df in zip([2, 3], [a1_res, a2_res]):
            ws = sh.get_worksheet(ws_idx); ws.batch_clear(["A2:G20"])
            ws.update(range_name='A3', values=df.values.tolist())

def process_jing_tao(files):
    df = None
    for f in files:
        f.seek(0); is_excel = f.name.lower().endswith(('.xlsx', '.xls'))
        if is_excel:
            try:
                xls = pd.ExcelFile(f); ts = next((s for s in xls.sheet_names if '靜桃' in s), xls.sheet_names[0])
                df_t = pd.read_excel(xls, sheet_name=ts, header=None, nrows=50)
                for idx, row in df_t.iterrows():
                    if '通報日期' in " ".join([str(x) for x in row if pd.notna(x)]): f.seek(0); df = pd.read_excel(xls, sheet_name=ts, skiprows=idx); break
                if df is not None: break
            except: pass
        if df is None:
            f.seek(0); rb = f.read()
            for enc in ['utf-8-sig', 'utf-8', 'cp950', 'big5']:
                try:
                    text = rb.decode(enc, errors='ignore'); lines = text.splitlines()
                    for idx, line in enumerate(lines[:50]):
                        if '通報日期' in line: f.seek(0); df = pd.read_csv(f, encoding=enc, skiprows=idx, engine='python', on_bad_lines='skip'); break
                    if df is not None: break
                except: continue
        if df is not None: break
    if df is None: return
    df.columns = [str(c).strip().replace('\u3000', '').replace('\n', '') for c in df.columns]
    dc = next((c for c in df.columns if '通報日期' in c), None); uc = next((c for c in df.columns if '所別' in c or ('單位' in c and '舉發單位' not in c)), None)
    c22 = next((c for c in df.columns if re.search(r'22.{0,3}0?6|夜間|深夜', c)), None); c06 = next((c for c in df.columns if re.search(r'0?6.{0,3}22|日間|白天', c)), None)
    def prd(v):
        try: p = re.split(r'[/\-]', str(v).strip().split(' ')[0]); return pd.Timestamp(year=int(p[0])+1911, month=int(p[1]), day=int(p[2]))
        except: return pd.NaT
    df['_date'] = df[dc].apply(prd); td = datetime.now(); ed = td - timedelta(days=1); sd = ed - timedelta(days=6)
    df_p = df[(df['_date'] >= sd) & (df['_date'] <= ed)]; vd = df['_date'].dropna()
    cs = f"({vd.min().year-1911}{vd.min().strftime('%m%d')}-{ed.year-1911}{ed.strftime('%m%d')})" if not vd.empty else ""
    res = []; tp22 = tp06 = ta22 = ta06 = tt = 0
    for kw, nm in zip(['聖亭', '龍潭', '中興', '石門', '高平', '三和', '警備', '交通'], ['聖亭所', '龍潭所', '中興所', '石門所', '高平所', '三和所', '警備隊', '交通分隊']):
        ma = df[uc].astype(str).str.contains(kw, na=False); mp = df_p[uc].astype(str).str.contains(kw, na=False)
        p22 = df_p[mp][c22].astype(str).str.contains(r'^V$', regex=True, na=False).sum() if c22 else 0
        p06 = df_p[mp][c06].astype(str).str.contains(r'^V$', regex=True, na=False).sum() if c06 else 0
        a22 = df[ma][c22].astype(str).str.contains(r'^V$', regex=True, na=False).sum() if c22 else 0
        a06 = df[ma][c06].astype(str).str.contains(r'^V$', regex=True, na=False).sum() if c06 else 0
        tot = len(df[ma]) if not c22 and not c06 else a22 + a06
        res.append([nm, p22, p06, a22, a06, tot]); tp22 += p22; tp06 += p06; ta22 += a22; ta06 += a06; tt += tot
    res.insert(0, ['合計', tp22, tp06, ta22, ta06, tt])
    df_res = pd.DataFrame(res, columns=pd.MultiIndex.from_arrays([['統計期間', f'本期({sd.strftime("%m%d")}-{ed.strftime("%m%d")})', f'本期({sd.strftime("%m%d")}-{ed.strftime("%m%d")})', f'累計{cs}', f'累計{cs}', '總計'], ['', c22 or "22-6時", c06 or "6-22時", c22 or "22-6時", c06 or "6-22時", '']]))
    st.write("📊 **「靜桃計畫」大執法專案統計表：**"); st.dataframe(df_res, use_container_width=True)
    if GCP_CREDS:
        gc = gspread.service_account_from_dict(GCP_CREDS); sh = gc.open_by_url(GOOGLE_SHEET_URL)
        ws = sh.worksheet("靜桃計畫") if "靜桃計畫" in [s.title for s in sh.worksheets()] else sh.add_worksheet(title="靜桃計畫", rows="30", cols="10")
        ws.clear(); ws.update(range_name='A1', values=[['「靜桃計畫」大執法專案統計表'], [t[0] for t in df_res.columns], [t[1] for t in df_res.columns]] + df_res.values.tolist())

# ==========================================
# 4. 首頁與分流器
# ==========================================
with st.sidebar:
    st.title("🚓 龍潭分局戰情室")
    app_mode = st.selectbox("功能模組", ["🏠 智慧批次處理中心", "📂 PDF 轉 PPTX 工具"])

if app_mode == "🏠 智慧批次處理中心":
    st.header("📈 交通數據全自動批次處理中心")
    uploads = st.file_uploader("📂 拖入所有報表檔案", type=["xlsx", "csv", "xls"], accept_multiple_files=True)
    if uploads:
        file_hash = sum([f.size for f in uploads]) + len(uploads)
        if st.session_state.get("last_processed_hash") != file_hash:
            cat_files = {"科技執法": [], "重大違規": [], "超載統計": [], "強化專案": [], "交通事故": [], "靜桃計畫": []}
            for f in uploads:
                name = f.name.lower()
                if any(k in name for k in ["list", "地點", "科技"]): cat_files["科技執法"].append(f)
                elif any(k in name for k in ["stone", "超載"]): cat_files["超載統計"].append(f)
                elif any(k in name for k in ["重大", "重點"]): cat_files["重大違規"].append(f)
                elif any(k in name for k in ["強化", "專案", "砂石", "大貨", "r17", "法條", "自選匯出"]): cat_files["強化專案"].append(f)
                elif any(k in name for k in ["a1", "a2", "事故", "案件統計"]): cat_files["交通事故"].append(f)
                elif any(k in name for k in ["靜桃", "噪音", "改裝車", "總表", "詳細資料"]): cat_files["靜桃計畫"].append(f)
            try:
                if cat_files["科技執法"]: with st.status("📸 處理【科技執法】..."): process_tech_enforcement(cat_files["科技執法"])
                if cat_files["超載統計"]: with st.status("🚛 處理【超載統計】..."): process_overload(cat_files["超載統計"])
                if cat_files["重大違規"]: with st.status("🚨 處理【重大交通違規】..."): process_major(cat_files["重大違規"])
                if cat_files["強化專案"]: with st.status("🔥 處理【強化專案】..."): process_project(cat_files["強化專案"])
                if cat_files["交通事故"]: with st.status("🚑 處理【交通事故】..."): process_accident(cat_files["交通事故"])
                if cat_files["靜桃計畫"]: with st.status("🤫 處理【靜桃計畫】..."): process_jing_tao(cat_files["靜桃計畫"])
                st.session_state["last_processed_hash"] = file_hash; st.balloons()
            except Exception as e: st.error(f"⚠️ 錯誤：{e}"); st.write(traceback.format_exc())

elif app_mode == "📂 PDF 轉 PPTX 工具":
    st.header("📂 PDF 行政文書轉 PPTX 簡報")
    pdf_file = st.file_uploader("上傳 PDF 檔案", type=["pdf"])
    if pdf_file and st.button("🚀 開始轉換"):
        with st.spinner("轉換中..."):
            try:
                images = convert_from_bytes(pdf_file.read(), dpi=200); prs = Presentation()
                for img in images:
                    slide = prs.slides.add_slide(prs.slide_layouts[6]); img_io = io.BytesIO()
                    img.save(img_io, format='PNG'); slide.shapes.add_picture(io.BytesIO(img_io.getvalue()), 0, 0, width=prs.slide_width, height=prs.slide_height)
                pptx_io = io.BytesIO(); prs.save(pptx_io); st.success("✅ 轉換完成！")
                st.download_button("📥 下載 PPTX", data=pptx_io.getvalue(), file_name=f"{pdf_file.name.replace('.pdf', '')}_轉換.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            except Exception as e: st.error(f"錯誤：{e}")
